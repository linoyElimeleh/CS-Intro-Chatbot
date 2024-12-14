import logging
import certifi
import os
from pinecone import Pinecone
import pinecone

os.environ["SSL_CERT_FILE"] = certifi.where()

from dotenv import load_dotenv

load_dotenv()
from langchain_community.document_loaders import (
    UnstructuredPowerPointLoader,
    TextLoader,
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from pptx import Presentation
from docx import Document as DocxDocument

from langchain_openai import OpenAIEmbeddings
from langchain_pinecone import PineconeVectorStore

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
embeddings = OpenAIEmbeddings(model="text-embedding-ada-002")


# Translate the data into Pinecone
def ingest_docs():
    try:
        clear_pinecone_vector_store()
    except Exception as e:
        logger.error(f"Error during document ingestion: {str(e)}")

    base_path = "intro-to-cs-public"
    documents = []

    logger.info(f"Start ingesting docs from {base_path}")

    for root, dirs, files in os.walk(base_path):
        for file in files:
            if file.endswith((".pptx", ".java", ".docx")):
                file_path = os.path.join(root, file)
                try:
                    if file.endswith(".pptx"):
                        prs = Presentation(file_path)
                        for i, slide in enumerate(prs.slides):
                            slide_content = ""
                            for shape in slide.shapes:
                                if hasattr(shape, 'text'):
                                    slide_content += shape.text + "\n"
                            if slide_content.strip():
                                documents.append(Document(
                                    page_content=slide_content.strip(),
                                    metadata={"source": f"{file_path}:Slide{i+1}"}
                                ))
                                logger.info(f"Processed {file_path}:Slide{i+1}")
                    
                    elif file.endswith(".docx"):
                        doc = DocxDocument(file_path)
                        doc_content = ""
                        for para in doc.paragraphs:
                            doc_content += para.text + "\n"
                        if doc_content.strip():
                            documents.append(Document(
                                page_content=doc_content.strip(),
                                metadata={"source": f"{file_path}"}
                            ))
                            logger.info(f"Processed {file_path}")

                    # Java file handling (unchanged)
                    else:
                        loader = TextLoader(file_path)
                        content = loader.load()[0].page_content
                        java_splitter = RecursiveCharacterTextSplitter(
                            separators=["\n\n"],
                            chunk_size=1000,
                            chunk_overlap=0
                        )
                        java_docs = java_splitter.create_documents([content])
                        for j, doc in enumerate(java_docs):
                            doc.metadata["source"] = f"{file_path}:Section{j+1}"
                        documents.extend(java_docs)
                        logger.info(f"Processed Java file {file_path} into {len(java_docs)} sections")

                except Exception as e:
                    logger.error(f"Error processing {file_path}: {str(e)}")

    logger.info(f"Total processed documents: {len(documents)}")

    for i, doc in enumerate(documents):
        if not hasattr(doc, 'metadata') or 'source' not in doc.metadata:
            doc.metadata['source'] = f"document_{i}"

    logger.info(f"Going to add {len(documents)} documents to Pinecone")
    PineconeVectorStore.from_documents(
        documents, embeddings, index_name=os.getenv("PINECONE_INDEX_NAME")
    )
    logger.info("****Loading to vectorstore done ***")


def clear_pinecone_vector_store():
    try:
        # Initialize Pinecone with API key and environment
        pc = Pinecone(
            api_key=os.getenv("PINECONE_API_KEY")
        )
        # Initialize the index with the name and no need for explicit host
        index_name = os.getenv("PINECONE_INDEX_NAME")
        
        index = pc.Index(index_name)

        # Delete all vectors in the index
        index.delete(delete_all=True)
        logger.info("Pinecone vector store cleared successfully.")
    
    except pinecone.exceptions.NotFoundException:
        logger.warning("Pinecone index not found. It may have already been deleted or not yet created.")
    
    except Exception as e:
        logger.error(f"Error clearing Pinecone vector store: {str(e)}")


if __name__ == "__main__":
    ingest_docs()
